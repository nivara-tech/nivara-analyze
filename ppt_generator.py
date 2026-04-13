"""
PowerPoint Generator for Eureka Forbes P&L Analysis.
Creates a professional PPT with multiple analysis slides.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import os

# Color palette
DARK_BLUE = RGBColor(0x1F, 0x4E, 0x79)
MEDIUM_BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
GREEN = RGBColor(0x00, 0x70, 0x30)
RED = RGBColor(0xC0, 0x00, 0x00)
DARK_RED = RGBColor(0x8B, 0x00, 0x00)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x80, 0x80, 0x80)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

FY_MONTHS = ["Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec", "Jan", "Feb", "Mar"]


def _n(val, default=0):
    """Coerce None to a numeric default."""
    return val if val is not None else default


def _var_pct(a, b):
    """Safe variance percentage."""
    a, b = _n(a), _n(b)
    return ((a - b) / abs(b) * 100) if b != 0 else 0


def _add_slide(prs, layout_idx=6):
    """Add a blank slide with Nivara Analyze branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    # Bottom-right branding
    txBox = slide.shapes.add_textbox(Inches(10.3), Inches(7.1), Inches(2.8), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Powered by Nivara Analyze"
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.font.italic = True
    p.alignment = PP_ALIGN.RIGHT
    return slide


def _add_title_bar(slide, title_text, subtitle_text=""):
    """Add a dark blue title bar at the top of a slide."""
    from pptx.util import Inches, Pt
    # Title bar background
    left, top, width, height = Inches(0), Inches(0), Inches(13.33), Inches(0.8)
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title text
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)


def _add_text_box(slide, left, top, width, height, text, font_size=10,
                  bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def _add_table(slide, left, top, width, height, rows, cols):
    """Add a table and return it."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table


def _style_header_row(table, col_count):
    """Style header row of a table."""
    for ci in range(col_count):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = WHITE
            paragraph.font.size = Pt(9)
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _style_data_cell(cell, value, is_currency=True, is_total=False, highlight_color=None):
    """Style a data cell."""
    if isinstance(value, (int, float)):
        if is_currency:
            cell.text = f"{value:,.1f}"
        else:
            cell.text = f"{value:.1f}%"
    else:
        cell.text = str(value) if value else "-"

    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(8)
        p.alignment = PP_ALIGN.RIGHT
        if is_total:
            p.font.bold = True
        if highlight_color:
            p.font.color.rgb = highlight_color

    if is_total:
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BLUE

    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _color_for_variance(pct, is_cost=False):
    """Return color based on variance. For costs, positive = bad."""
    if is_cost:
        pct = -pct
    if pct > 5:
        return GREEN
    elif pct < -5:
        return RED
    return BLACK


def _add_kpi_card(slide, left, top, width, height, title, value, subtitle="", color=DARK_BLUE):
    """Add a KPI card with title, big value, and subtitle."""
    # Card background
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)

    # Title
    _add_text_box(slide, left + Inches(0.1), top + Inches(0.05),
                  width - Inches(0.2), Inches(0.25),
                  title, font_size=8, color=GRAY)

    # Value
    _add_text_box(slide, left + Inches(0.1), top + Inches(0.28),
                  width - Inches(0.2), Inches(0.35),
                  value, font_size=18, bold=True, color=color)

    # Subtitle
    if subtitle:
        sub_color = GREEN if "+" in subtitle or "Beat" in subtitle else RED if "-" in subtitle or "Miss" in subtitle else GRAY
        _add_text_box(slide, left + Inches(0.1), top + Inches(0.6),
                      width - Inches(0.2), Inches(0.2),
                      subtitle, font_size=8, color=sub_color)


def _add_narrative_slide(prs, title, subtitle, bullets, bullet_color=BLACK):
    """Add a slide with narrative bullet points (from LLM insights)."""
    slide = _add_slide(prs)
    _add_title_bar(slide, title, subtitle)

    if isinstance(bullets, str):
        bullets = [bullets]

    y = 1.3
    for bullet in bullets:
        if not bullet:
            continue
        # Auto-color: green for positive, red for negative
        color = bullet_color
        if any(w in bullet.lower() for w in ["beat", "exceeded", "strong", "growth", "improved", "positive"]):
            color = GREEN
        elif any(w in bullet.lower() for w in ["miss", "decline", "risk", "concern", "below", "negative", "increased cost"]):
            color = RED

        _add_text_box(slide, Inches(0.7), Inches(y), Inches(11.5), Inches(0.5),
                      f"•  {bullet}", font_size=11, color=color)
        y += 0.5
        if y > 6.8:
            break


def generate_ppt(analysis, output_path, insights=None):
    """Generate the full analysis PPT. If insights dict is provided, adds LLM narrative slides."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    a = analysis
    is_full_year = a.review_month == "Mar"
    period_label = f"Full Year {a.review_fy}" if is_full_year else f"{a.review_month} {a.review_fy}"

    # ==================== SLIDE 1: Title Slide ====================
    slide = _add_slide(prs)
    # Big logo area
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    _add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                  "EUREKA FORBES LIMITED", font_size=36, bold=True, color=WHITE,
                  alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                  f"P&L Review — {period_label}", font_size=24, color=RGBColor(0xBB, 0xCC, 0xDD),
                  alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
                  "Monthly Business Review | Confidential", font_size=14, color=RGBColor(0x88, 0x99, 0xAA),
                  alignment=PP_ALIGN.CENTER)
    yr_text = "Full Year Results" if is_full_year else "Monthly Review"
    _add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                  yr_text, font_size=16, bold=True, color=WHITE,
                  alignment=PP_ALIGN.CENTER)

    # ==================== SLIDE 2: BIG PICTURE ====================
    slide = _add_slide(prs)
    _add_title_bar(slide, f"The Big Picture — {period_label}")

    # Gather all big picture metrics for current month, AOP, LY, and YTD
    bp_metrics = {
        "Total Gross Sales": {"label": "Gross Sales"},
        "Total Net Sales": {"label": "Net Sales"},
        "Gross Margin": {"label": "Gross Margin"},
        "Gross Margin %": {"label": "Gross Margin %", "is_pct": True},
        "EBITDA (post allocation)": {"label": "EBITDA"},
        "EBITDA %": {"label": "EBITDA Margin", "is_pct": True},
        "Profit After Tax": {"label": "PAT"},
        "Advent Adjusted EBITDA": {"label": "Advent Adj. EBITDA"},
        "Advent Adjusted EBITDA %": {"label": "Advent Adj. EBITDA %", "is_pct": True},
    }

    # Big KPI cards row 1: Gross Sales, Net Sales, EBITDA, EBITDA Margin
    def _bp_val(key):
        return _n(a.current_month.get(key))

    def _bp_aop(key):
        return _n(a.aop_month.get(key))

    def _bp_ly(key):
        return _n(a.ly_same_month.get(key))

    def _bp_ytd(key):
        return _n(a.ytd_actual.get(key))

    def _bp_ytd_aop(key):
        return _n(a.ytd_aop.get(key))

    # Row 1: Big numbers
    row1 = [
        ("Gross Sales", "Total Gross Sales", False),
        ("Net Sales", "Total Net Sales", False),
        ("Gross Margin", "Gross Margin", False),
        ("EBITDA", "EBITDA (post allocation)", False),
        ("Advent Adj. EBITDA", "Advent Adjusted EBITDA", False),
    ]

    for i, (title, key, is_pct) in enumerate(row1):
        val = _bp_val(key)
        aop = _bp_aop(key)
        ly = _bp_ly(key)
        ach = (val / aop * 100) if aop else 0
        yoy = _var_pct(val, ly)
        color = GREEN if ach >= 100 else ORANGE if ach >= 90 else RED
        _add_kpi_card(slide, Inches(0.3 + i * 2.55), Inches(1.1), Inches(2.4), Inches(1.1),
                      title,
                      f"Rs {val:,.0f} Cr",
                      f"AOP {ach:.0f}% | YoY {yoy:+.1f}%",
                      color)

    # Row 2: Margins + Growth
    ns_val = _bp_val("Total Net Sales")
    ns_ly = _bp_ly("Total Net Sales")
    ns_yoy = _var_pct(ns_val, ns_ly)

    gs_val = _bp_val("Total Gross Sales")
    gs_ly = _bp_ly("Total Gross Sales")
    gs_yoy = _var_pct(gs_val, gs_ly)

    ebitda_margin = _bp_val("EBITDA %")
    ebitda_margin_ly = _bp_ly("EBITDA %")

    gm_pct_val = _bp_val("Gross Margin %")
    gm_pct_ly = _bp_ly("Gross Margin %")

    advent_margin = _bp_val("Advent Adjusted EBITDA %")
    advent_margin_ly = _bp_ly("Advent Adjusted EBITDA %")

    row2 = [
        ("Net Sales Growth (YoY)", f"{ns_yoy:+.1f}%", f"Rs {ns_val:,.0f} Cr vs LY Rs {ns_ly:,.0f} Cr",
         GREEN if ns_yoy > 0 else RED),
        ("Gross Sales Growth (YoY)", f"{gs_yoy:+.1f}%", f"Rs {gs_val:,.0f} Cr vs LY Rs {gs_ly:,.0f} Cr",
         GREEN if gs_yoy > 0 else RED),
        ("Gross Margin %", f"{gm_pct_val:.1f}%", f"LY: {gm_pct_ly:.1f}% | Delta: {gm_pct_val - gm_pct_ly:+.1f}pp",
         GREEN if gm_pct_val >= gm_pct_ly else RED),
        ("EBITDA Margin", f"{ebitda_margin:.1f}%", f"LY: {ebitda_margin_ly:.1f}% | Delta: {ebitda_margin - ebitda_margin_ly:+.1f}pp",
         GREEN if ebitda_margin >= ebitda_margin_ly else RED),
        ("Advent Adj. EBITDA Margin", f"{advent_margin:.1f}%", f"LY: {advent_margin_ly:.1f}% | Delta: {advent_margin - advent_margin_ly:+.1f}pp",
         GREEN if advent_margin >= advent_margin_ly else RED),
    ]

    for i, (title, value, subtitle, color) in enumerate(row2):
        _add_kpi_card(slide, Inches(0.3 + i * 2.55), Inches(2.5), Inches(2.4), Inches(1.1),
                      title, value, subtitle, color)

    # Row 3: YTD / Full Year view
    _add_text_box(slide, Inches(0.3), Inches(4.0), Inches(12), Inches(0.4),
                  "YTD / FULL YEAR VIEW" if is_full_year else f"YTD (Apr–{a.review_month})",
                  font_size=13, bold=True, color=DARK_BLUE)

    ytd_items = [
        ("YTD Gross Sales", "Total Gross Sales"),
        ("YTD Net Sales", "Total Net Sales"),
        ("YTD EBITDA", "EBITDA (post allocation)"),
        ("YTD PAT", "Profit After Tax"),
        ("YTD Advent Adj. EBITDA", "Advent Adjusted EBITDA"),
    ]

    for i, (title, key) in enumerate(ytd_items):
        ytd_val = _bp_ytd(key)
        ytd_aop = _bp_ytd_aop(key)
        ach = (ytd_val / ytd_aop * 100) if ytd_aop else 0
        color = GREEN if ach >= 100 else ORANGE if ach >= 90 else RED
        _add_kpi_card(slide, Inches(0.3 + i * 2.55), Inches(4.5), Inches(2.4), Inches(1.0),
                      title,
                      f"Rs {ytd_val:,.0f} Cr",
                      f"AOP Achieved: {ach:.0f}%",
                      color)

    # Advent EBITDA Bridge at bottom
    _add_text_box(slide, Inches(0.3), Inches(5.8), Inches(6), Inches(0.35),
                  "ADVENT ADJUSTED EBITDA BRIDGE", font_size=11, bold=True, color=DARK_BLUE)

    bridge_items = [
        ("Reported EBITDA", "Reported EBITDA"),
        ("(+) Ind AS Adjustment", "Ind AS Adjustment"),
        ("(+) ESOP Cost Add-back", "ESOP Cost Add-back"),
        ("(+) One-time Items", "One-time Items"),
        ("= Advent Adj. EBITDA", "Advent Adjusted EBITDA"),
    ]

    cols_br = ["Item", f"{a.review_month} {a.review_fy}", "AOP", "Var"]
    table = _add_table(slide, Inches(0.3), Inches(6.15), Inches(8), Inches(0.3 + len(bridge_items) * 0.22),
                       len(bridge_items) + 1, len(cols_br))

    table.columns[0].width = Inches(3.0)
    for ci in range(1, len(cols_br)):
        table.columns[ci].width = Inches(1.65)

    for ci, h in enumerate(cols_br):
        table.cell(0, ci).text = h
    _style_header_row(table, len(cols_br))

    for ri, (display, key) in enumerate(bridge_items):
        row_idx = ri + 1
        is_total = "=" in display

        cell = table.cell(row_idx, 0)
        cell.text = display
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.bold = is_total
        if is_total:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BLUE

        actual = _n(a.current_month.get(key))
        aop = _n(a.aop_month.get(key))
        var = actual - aop

        _style_data_cell(table.cell(row_idx, 1), actual, True, is_total)
        _style_data_cell(table.cell(row_idx, 2), aop, True, is_total)
        _style_data_cell(table.cell(row_idx, 3), var, True, is_total,
                         GREEN if var > 0 else RED if var < 0 else BLACK)

    # ==================== SLIDE 3: Executive Summary ====================
    slide = _add_slide(prs)
    _add_title_bar(slide, f"Executive Summary — {period_label}")

    # KPI Cards
    ns = _n(a.current_month.get("Total Net Sales"))
    ns_aop = _n(a.aop_month.get("Total Net Sales"))
    ns_ach = (ns / ns_aop * 100) if ns_aop else 0

    gm = _n(a.current_month.get("Gross Margin"))
    gm_pct = _n(a.current_month.get("Gross Margin %"))

    ebitda = _n(a.current_month.get("EBITDA (post allocation)"))
    ebitda_aop = _n(a.aop_month.get("EBITDA (post allocation)"))
    ebitda_ach = (ebitda / ebitda_aop * 100) if ebitda_aop else 0

    pat = _n(a.current_month.get("Profit After Tax"))
    pat_aop = _n(a.aop_month.get("Profit After Tax"))

    pbt = _n(a.current_month.get("Profit Before Tax"))
    pbt_pct = _n(a.current_month.get("PBT %"))

    pat_ach = (pat / pat_aop * 100) if pat_aop else 0

    cards = [
        ("Net Sales", f"Rs {ns:.0f} Cr", f"Achieved AOP {ns_ach:.0f}%",
         GREEN if ns_ach >= 100 else ORANGE if ns_ach >= 90 else RED),
        ("Gross Margin", f"{gm_pct:.1f}%", f"Rs {gm:.0f} Cr", DARK_BLUE),
        ("EBITDA", f"Rs {ebitda:.0f} Cr", f"Achieved AOP {ebitda_ach:.0f}%",
         GREEN if ebitda_ach >= 100 else ORANGE if ebitda_ach >= 90 else RED),
        ("PBT", f"Rs {pbt:.0f} Cr", f"Margin: {pbt_pct:.1f}%", DARK_BLUE),
        ("PAT", f"Rs {pat:.0f} Cr", f"Achieved AOP {pat_ach:.0f}%",
         GREEN if pat_ach >= 100 else ORANGE if pat_ach >= 90 else RED),
    ]

    for i, (title, value, subtitle, color) in enumerate(cards):
        _add_kpi_card(slide, Inches(0.5 + i * 2.5), Inches(1.1), Inches(2.3), Inches(0.9),
                      title, value, subtitle, color)

    # Key Highlights
    _add_text_box(slide, Inches(0.5), Inches(2.3), Inches(6), Inches(0.4),
                  "KEY HIGHLIGHTS", font_size=12, bold=True, color=DARK_BLUE)

    for i, h in enumerate(a.highlights[:8]):
        icon = "+" if any(w in h for w in ["BEAT", "above", "Growth"]) else "-" if any(w in h for w in ["MISS", "below"]) else ">"
        color = GREEN if icon == "+" else RED if icon == "-" else BLACK
        _add_text_box(slide, Inches(0.7), Inches(2.8 + i * 0.35), Inches(12), Inches(0.35),
                      f"  {icon}  {h}", font_size=9, color=color)

    # ==================== SLIDE 2B: Product Segment Growth (EWP, VC, etc.) ====================
    slide = _add_slide(prs)
    _add_title_bar(slide, f"Product Segment Performance — {a.review_month} {a.review_fy}",
                   "EWP & VC Growth | Segment-wise Gross Sales")

    product_lines = [
        "Gross Sales - EWP", "Gross Sales - VC", "Gross Sales - NEWP",
        "Gross Sales - Air Purifier", "Gross Sales - Softeners", "Gross Sales - Others",
        "Total Product Gross Sales",
    ]

    cols_seg = ["Product Segment", f"{a.review_month} Actual", f"{a.review_month} AOP",
                "Ach %", "Prev Month", "MoM %", f"LY {a.review_month}", "YoY Growth %"]
    num_seg_rows = len(product_lines) + 1
    table = _add_table(slide, Inches(0.3), Inches(1.2), Inches(12.7), Inches(0.4 + len(product_lines) * 0.55),
                       num_seg_rows, len(cols_seg))

    table.columns[0].width = Inches(3.0)
    for ci in range(1, len(cols_seg)):
        table.columns[ci].width = Inches(1.4)

    for ci, h in enumerate(cols_seg):
        table.cell(0, ci).text = h
    _style_header_row(table, len(cols_seg))

    for ri, label in enumerate(product_lines):
        row_idx = ri + 1
        is_total = "Total" in label

        cell = table.cell(row_idx, 0)
        display = label.replace("Gross Sales - ", "")
        cell.text = display
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = is_total
        if is_total:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BLUE

        actual = _n(a.current_month.get(label))
        aop = _n(a.aop_month.get(label))
        prev = _n(a.prev_month.get(label))
        ly = _n(a.ly_same_month.get(label))

        ach_pct = (actual / aop * 100) if aop else 0
        mom_pct = _var_pct(actual, prev)
        yoy_pct = _var_pct(actual, ly)

        _style_data_cell(table.cell(row_idx, 1), actual, True, is_total)
        _style_data_cell(table.cell(row_idx, 2), aop, True, is_total)
        _style_data_cell(table.cell(row_idx, 3), ach_pct, False, is_total,
                         GREEN if ach_pct >= 100 else RED if ach_pct < 90 else ORANGE)
        _style_data_cell(table.cell(row_idx, 4), prev, True, is_total)
        _style_data_cell(table.cell(row_idx, 5), mom_pct, False, is_total,
                         _color_for_variance(mom_pct, False))
        _style_data_cell(table.cell(row_idx, 6), ly, True, is_total)
        _style_data_cell(table.cell(row_idx, 7), yoy_pct, False, is_total,
                         _color_for_variance(yoy_pct, False))

    # EWP+VC combined growth KPIs
    ewp_actual = _n(a.current_month.get("Gross Sales - EWP"))
    ewp_ly = _n(a.ly_same_month.get("Gross Sales - EWP"))
    vc_actual = _n(a.current_month.get("Gross Sales - VC"))
    vc_ly = _n(a.ly_same_month.get("Gross Sales - VC"))

    ewp_yoy = _var_pct(ewp_actual, ewp_ly)
    vc_yoy = _var_pct(vc_actual, vc_ly)
    combined = ewp_actual + vc_actual
    combined_ly = ewp_ly + vc_ly
    combined_yoy = _var_pct(combined, combined_ly)

    y_kpi = 1.2 + 0.4 + len(product_lines) * 0.55 + 0.4
    _add_kpi_card(slide, Inches(0.5), Inches(y_kpi), Inches(3.5), Inches(0.9),
                  "EWP YoY Growth", f"{ewp_yoy:+.1f}%",
                  f"Rs {ewp_actual:.0f} Cr vs LY Rs {ewp_ly:.0f} Cr",
                  GREEN if ewp_yoy > 0 else RED)
    _add_kpi_card(slide, Inches(4.5), Inches(y_kpi), Inches(3.5), Inches(0.9),
                  "VC YoY Growth", f"{vc_yoy:+.1f}%",
                  f"Rs {vc_actual:.0f} Cr vs LY Rs {vc_ly:.0f} Cr",
                  GREEN if vc_yoy > 0 else RED)
    _add_kpi_card(slide, Inches(8.5), Inches(y_kpi), Inches(3.5), Inches(0.9),
                  "EWP+VC Combined YoY", f"{combined_yoy:+.1f}%",
                  f"Rs {combined:.0f} Cr vs LY Rs {combined_ly:.0f} Cr",
                  GREEN if combined_yoy > 0 else RED)

    # ==================== SLIDE 3: P&L Overview - Monthly Comparison ====================
    slide = _add_slide(prs)
    _add_title_bar(slide, f"P&L Overview — {a.review_month} {a.review_fy}",
                   "Current Month vs AOP | Previous Month | Last Year Same Month")

    key_pnl_items = [
        ("Total Net Sales", False),
        ("Total COGS", True),
        ("Gross Margin", False),
        ("Gross Margin %", False),
        ("Total Employee Related Costs", True),
        ("Total Adv & Sales Promo", True),
        ("Total Service Charges", True),
        ("Total Other Direct Costs", True),
        ("Total Direct Costs", True),
        ("Contribution Profit (GM-DC)", False),
        ("Contribution Profit %", False),
        ("Total Indirect Costs (A+B)", True),
        ("EBITDA Operating", False),
        ("EBITDA Operating %", False),
        ("Common + Mfg Overheads", True),
        ("EBITDA (post allocation)", False),
        ("EBITDA %", False),
        ("Depreciation", True),
        ("EBIT", False),
        ("Interest Cost", True),
        ("Profit Before Tax", False),
        ("PBT %", False),
        ("Tax", True),
        ("Profit After Tax", False),
        ("PAT %", False),
    ]

    cols = ["Line Item", f"{a.review_month} Actual", f"{a.review_month} AOP",
            "Var %", "Prev Month", "MoM %", f"LY {a.review_month}", "YoY %"]
    num_rows = len(key_pnl_items) + 1
    table = _add_table(slide, Inches(0.3), Inches(1.0), Inches(12.7), Inches(6.2),
                       num_rows, len(cols))

    # Set column widths
    table.columns[0].width = Inches(3.0)
    for ci in range(1, len(cols)):
        table.columns[ci].width = Inches(1.4)

    # Headers
    for ci, h in enumerate(cols):
        table.cell(0, ci).text = h
    _style_header_row(table, len(cols))

    # Data
    for ri, (label, is_cost) in enumerate(key_pnl_items):
        row_idx = ri + 1
        is_pct = "%" in label
        is_total = any(kw in label for kw in ["Total", "Gross Margin", "Contribution", "EBITDA", "EBIT", "Profit", "PAT"])

        # Line item name
        cell = table.cell(row_idx, 0)
        cell.text = f"  {label}" if not is_total else label
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.bold = is_total
        if is_total:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BLUE

        actual = a.current_month.get(label, 0)
        aop = a.aop_month.get(label, 0)
        prev = a.prev_month.get(label, 0)
        ly = a.ly_same_month.get(label, 0)

        actual, aop, prev, ly = _n(actual), _n(aop), _n(prev), _n(ly)
        var_pct = _var_pct(actual, aop)
        mom_pct = _var_pct(actual, prev)
        yoy_pct = _var_pct(actual, ly)

        _style_data_cell(table.cell(row_idx, 1), actual, not is_pct, is_total)
        _style_data_cell(table.cell(row_idx, 2), aop, not is_pct, is_total)
        _style_data_cell(table.cell(row_idx, 3), var_pct, False, is_total,
                         _color_for_variance(var_pct, is_cost))
        _style_data_cell(table.cell(row_idx, 4), prev, not is_pct, is_total)
        _style_data_cell(table.cell(row_idx, 5), mom_pct, False, is_total,
                         _color_for_variance(mom_pct, is_cost))
        _style_data_cell(table.cell(row_idx, 6), ly, not is_pct, is_total)
        _style_data_cell(table.cell(row_idx, 7), yoy_pct, False, is_total,
                         _color_for_variance(yoy_pct, is_cost))

    # ==================== SLIDE 4: Budget Achievement ====================
    slide = _add_slide(prs)
    _add_title_bar(slide, f"Budget Achievement (AOP) — {period_label}",
                   "Actual vs Annual Operating Plan")

    ba_items = list(a.budget_achievement.items())
    if ba_items:
        cols_ba = ["Metric", "Actual (Rs Cr)", "AOP (Rs Cr)", "Variance (Rs Cr)", "Achievement %", "Status"]
        table = _add_table(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.4 + len(ba_items) * 0.45),
                           len(ba_items) + 1, len(cols_ba))

        table.columns[0].width = Inches(3.5)
        for ci in range(1, len(cols_ba)):
            table.columns[ci].width = Inches(1.7)

        for ci, h in enumerate(cols_ba):
            table.cell(0, ci).text = h
        _style_header_row(table, len(cols_ba))

        for ri, (label, data) in enumerate(ba_items):
            row_idx = ri + 1
            cell = table.cell(row_idx, 0)
            cell.text = label
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.bold = True

            _style_data_cell(table.cell(row_idx, 1), data["actual"])
            _style_data_cell(table.cell(row_idx, 2), data["budget"])
            var_color = GREEN if data["variance"] > 0 else RED
            is_cost_item = any(kw in label for kw in ["Cost", "COGS", "Freight"])
            if is_cost_item:
                var_color = RED if data["variance"] > 0 else GREEN
            _style_data_cell(table.cell(row_idx, 3), data["variance"], True, False, var_color)
            _style_data_cell(table.cell(row_idx, 4), data["achievement_pct"], False)

            status = "BEAT" if data["achievement_pct"] >= 100 else "MISS"
            if is_cost_item:
                status = "MISS" if data["achievement_pct"] > 100 else "ON TRACK"
            cell = table.cell(row_idx, 5)
            cell.text = status
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.bold = True
                p.font.color.rgb = GREEN if status in ("BEAT", "ON TRACK") else RED
                p.alignment = PP_ALIGN.CENTER

    # YTD section below
    _add_text_box(slide, Inches(0.5), Inches(1.2 + 0.4 + len(ba_items) * 0.45 + 0.3),
                  Inches(12), Inches(0.4),
                  f"YTD Achievement (Apr-{a.review_month})", font_size=12, bold=True, color=DARK_BLUE)

    ytd_metrics = ["Total Net Sales", "EBITDA (post allocation)", "Profit After Tax"]
    y_start = 1.2 + 0.4 + len(ba_items) * 0.45 + 0.8
    for i, label in enumerate(ytd_metrics):
        ytd_a = a.ytd_actual.get(label, 0)
        ytd_b = a.ytd_aop.get(label, 0)
        ach = (ytd_a / ytd_b * 100) if ytd_b else 0
        color = GREEN if ach >= 100 else RED
        _add_kpi_card(slide, Inches(0.5 + i * 4), Inches(y_start), Inches(3.5), Inches(0.9),
                      f"YTD {label}",
                      f"Rs {ytd_a:.0f} Cr",
                      f"AOP: Rs {ytd_b:.0f} Cr | {ach:.1f}% Achieved",
                      color)

    # ==================== SLIDE 5: Quarterly Comparison ====================
    slide = _add_slide(prs)
    from analyzer import _get_quarter, _prev_quarter
    curr_q = _get_quarter(a.review_month)
    prev_q = _prev_quarter(curr_q)
    prev_q_fy = a.review_fy if curr_q != "Q1" else f"FY{a.review_fy_int % 100 - 1}"
    ly_q_fy = f"FY{a.review_fy_int % 100 - 1}"

    _add_title_bar(slide, f"Quarterly Comparison — {curr_q} {a.review_fy}",
                   f"{curr_q} {a.review_fy} vs {prev_q} {prev_q_fy} vs {curr_q} {ly_q_fy}")

    q_items = [
        "Total Net Sales", "Total COGS", "Gross Margin",
        "Total Direct Costs", "Contribution Profit (GM-DC)",
        "Total Indirect Costs (A+B)", "EBITDA (post allocation)",
        "Profit Before Tax", "Profit After Tax",
    ]

    cols_q = ["Line Item", f"{curr_q} {a.review_fy}", f"{curr_q} AOP",
              "Var %", f"{prev_q} Actual", "QoQ %",
              f"{curr_q} {ly_q_fy}", "YoY %"]
    table = _add_table(slide, Inches(0.3), Inches(1.1), Inches(12.7), Inches(0.4 + len(q_items) * 0.5),
                       len(q_items) + 1, len(cols_q))

    table.columns[0].width = Inches(3.0)
    for ci in range(1, len(cols_q)):
        table.columns[ci].width = Inches(1.4)

    for ci, h in enumerate(cols_q):
        table.cell(0, ci).text = h
    _style_header_row(table, len(cols_q))

    for ri, label in enumerate(q_items):
        row_idx = ri + 1
        is_cost = any(kw in label for kw in ["Cost", "COGS"])

        cell = table.cell(row_idx, 0)
        cell.text = label
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True

        cq = a.current_quarter.get(label, 0)
        aq = a.aop_quarter.get(label, 0)
        pq = a.prev_quarter.get(label, 0)
        lq = a.ly_same_quarter.get(label, 0)

        cq, aq, pq, lq = _n(cq), _n(aq), _n(pq), _n(lq)
        var_pct = _var_pct(cq, aq)
        qoq_pct = _var_pct(cq, pq)
        yoy_pct = _var_pct(cq, lq)

        _style_data_cell(table.cell(row_idx, 1), cq)
        _style_data_cell(table.cell(row_idx, 2), aq)
        _style_data_cell(table.cell(row_idx, 3), var_pct, False, False, _color_for_variance(var_pct, is_cost))
        _style_data_cell(table.cell(row_idx, 4), pq)
        _style_data_cell(table.cell(row_idx, 5), qoq_pct, False, False, _color_for_variance(qoq_pct, is_cost))
        _style_data_cell(table.cell(row_idx, 6), lq)
        _style_data_cell(table.cell(row_idx, 7), yoy_pct, False, False, _color_for_variance(yoy_pct, is_cost))

    # ==================== SLIDE 6: Full Year (if March) ====================
    if is_full_year:
        slide = _add_slide(prs)
        _add_title_bar(slide, f"Full Year Results — {a.review_fy}",
                       f"{a.review_fy} Actuals vs AOP vs Previous Year")

        fy_items = [
            "Total Net Sales", "Total COGS", "Gross Margin",
            "Total Employee Related Costs", "Total Adv & Sales Promo",
            "Total Service Charges", "Total Other Direct Costs",
            "Total Direct Costs", "Contribution Profit (GM-DC)",
            "Total Indirect Costs (A+B)", "EBITDA Operating",
            "Common + Mfg Overheads", "EBITDA (post allocation)",
            "Depreciation", "EBIT", "Profit Before Tax",
            "Tax", "Profit After Tax",
        ]

        cols_fy = ["Line Item", f"{a.review_fy} Actual", f"{a.review_fy} AOP",
                   "Var %", f"FY{a.review_fy_int % 100 - 1} Actual", "YoY %"]
        table = _add_table(slide, Inches(0.3), Inches(1.0), Inches(12.7), Inches(6.2),
                           len(fy_items) + 1, len(cols_fy))

        table.columns[0].width = Inches(3.5)
        for ci in range(1, len(cols_fy)):
            table.columns[ci].width = Inches(1.8)

        for ci, h in enumerate(cols_fy):
            table.cell(0, ci).text = h
        _style_header_row(table, len(cols_fy))

        for ri, label in enumerate(fy_items):
            row_idx = ri + 1
            is_cost = any(kw in label for kw in ["Cost", "COGS", "Depreciation", "Tax"])

            cell = table.cell(row_idx, 0)
            cell.text = label
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.bold = True

            fy_a = a.ytd_actual.get(label, 0)  # For March, YTD = full year
            fy_b = a.fy_aop.get(label, 0)
            ly = a.ly_full_year.get(label, 0)

            fy_a, fy_b, ly = _n(fy_a), _n(fy_b), _n(ly)
            var_pct = _var_pct(fy_a, fy_b)
            yoy_pct = _var_pct(fy_a, ly)

            _style_data_cell(table.cell(row_idx, 1), fy_a)
            _style_data_cell(table.cell(row_idx, 2), fy_b)
            _style_data_cell(table.cell(row_idx, 3), var_pct, False, False, _color_for_variance(var_pct, is_cost))
            _style_data_cell(table.cell(row_idx, 4), ly)
            _style_data_cell(table.cell(row_idx, 5), yoy_pct, False, False, _color_for_variance(yoy_pct, is_cost))

    # ==================== SLIDE 7: Revenue Trend Chart ====================
    slide = _add_slide(prs)
    _add_title_bar(slide, f"Revenue Trend — {a.review_fy}", "Monthly Net Sales with EBITDA Margin")

    chart_data = CategoryChartData()
    months_available = FY_MONTHS[:len(a.monthly_trend.get("Total Net Sales", []))]
    chart_data.categories = months_available

    ns_trend = a.monthly_trend.get("Total Net Sales", [])
    ebitda_trend = a.monthly_trend.get("EBITDA (post allocation)", [])
    aop_vals = []
    for m in months_available:
        k = f"{m}_FY{a.review_fy_int % 100}_AOP"
        # We need to get AOP from analysis... use fy_aop divided equally as approximation
        aop_vals.append(a.aop_month.get("Total Net Sales", 0))  # same month AOP

    chart_data.add_series("Net Sales (Actual)", ns_trend)
    chart_data.add_series("EBITDA", ebitda_trend)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.2), Inches(8), Inches(5.5),
        chart_data
    ).chart

    chart.has_legend = True
    chart.legend.include_in_layout = False

    # Side KPIs
    ns_total = sum(ns_trend)
    ebitda_total = sum(ebitda_trend)
    ebitda_margin = (ebitda_total / ns_total * 100) if ns_total else 0

    _add_kpi_card(slide, Inches(9), Inches(1.5), Inches(3.5), Inches(0.9),
                  f"YTD Net Sales", f"Rs {ns_total:.0f} Cr", "", DARK_BLUE)
    _add_kpi_card(slide, Inches(9), Inches(2.7), Inches(3.5), Inches(0.9),
                  f"YTD EBITDA", f"Rs {ebitda_total:.0f} Cr", f"Margin: {ebitda_margin:.1f}%", DARK_BLUE)

    # Best/worst month
    if ns_trend:
        best_idx = ns_trend.index(max(ns_trend))
        worst_idx = ns_trend.index(min(ns_trend))
        _add_kpi_card(slide, Inches(9), Inches(3.9), Inches(3.5), Inches(0.9),
                      "Best Month", f"{months_available[best_idx]}: Rs {ns_trend[best_idx]:.0f} Cr", "", GREEN)
        _add_kpi_card(slide, Inches(9), Inches(5.1), Inches(3.5), Inches(0.9),
                      "Weakest Month", f"{months_available[worst_idx]}: Rs {ns_trend[worst_idx]:.0f} Cr", "", RED)

    # ==================== SLIDE 8: Outliers - vs AOP ====================
    aop_outliers = [o for o in a.outliers if o.comparison_type == "vs_AOP"]
    if aop_outliers:
        slide = _add_slide(prs)
        _add_title_bar(slide, f"Outliers — Actual vs AOP Budget ({a.review_month} {a.review_fy})",
                       f"Items deviating >10% from Annual Operating Plan")

        _build_outlier_table(slide, aop_outliers[:18], "vs AOP")

    # ==================== SLIDE 9: Outliers - MoM ====================
    mom_outliers = [o for o in a.outliers if o.comparison_type == "MoM"]
    if mom_outliers:
        slide = _add_slide(prs)
        _add_title_bar(slide, f"Outliers — Month-over-Month ({a.review_month} {a.review_fy})",
                       "Items with >10% change vs previous month")

        _build_outlier_table(slide, mom_outliers[:18], "MoM")

    # ==================== SLIDE 10: Outliers - YoY ====================
    yoy_outliers = [o for o in a.outliers if o.comparison_type == "YoY"]
    if yoy_outliers:
        slide = _add_slide(prs)
        _add_title_bar(slide, f"Outliers — Year-over-Year ({a.review_month} {a.review_fy} vs {a.review_month} FY{a.review_fy_int % 100 - 1})",
                       "Items with >10% YoY change")

        _build_outlier_table(slide, yoy_outliers[:18], "YoY")

    # ==================== SLIDE 11: Outliers - QoQ ====================
    qoq_outliers = [o for o in a.outliers if o.comparison_type == "QoQ"]
    if qoq_outliers:
        slide = _add_slide(prs)
        _add_title_bar(slide, f"Outliers — Quarter-over-Quarter ({curr_q} {a.review_fy})",
                       "Items with >10% QoQ change (monthly avg basis)")

        _build_outlier_table(slide, qoq_outliers[:18], "QoQ")

    # ==================== SLIDE 12: Cost Deep Dive ====================
    slide = _add_slide(prs)
    _add_title_bar(slide, f"Cost Analysis Deep Dive — {a.review_month} {a.review_fy}")

    cost_items = [
        ("Employee Costs", "Total Employee Related Costs"),
        ("Adv & Sales Promo", "Total Adv & Sales Promo"),
        ("Service Charges", "Total Service Charges"),
        ("Other Direct", "Total Other Direct Costs"),
        ("Indirect Costs", "Total Indirect Costs (A+B)"),
        ("Overheads", "Common + Mfg Overheads"),
        ("Freight", "Freight"),
        ("COGS", "Total COGS"),
    ]

    cols_c = ["Cost Category", "Actual", "AOP", "Var %", "% to NS", "LY % to NS", "Delta"]
    table = _add_table(slide, Inches(0.3), Inches(1.1), Inches(12.7), Inches(0.4 + len(cost_items) * 0.5),
                       len(cost_items) + 1, len(cols_c))
    table.columns[0].width = Inches(2.8)
    for ci in range(1, len(cols_c)):
        table.columns[ci].width = Inches(1.6)

    for ci, h in enumerate(cols_c):
        table.cell(0, ci).text = h
    _style_header_row(table, len(cols_c))

    ns_actual = _n(a.current_month.get("Total Net Sales"), 1)
    ns_ly = _n(a.ly_same_month.get("Total Net Sales"), 1)

    for ri, (display_name, label) in enumerate(cost_items):
        row_idx = ri + 1
        actual = a.current_month.get(label, 0)
        aop = a.aop_month.get(label, 0)
        ly = a.ly_same_month.get(label, 0)

        actual, aop = _n(actual), _n(aop)
        var_pct = _var_pct(actual, aop)
        pct_ns = (actual / ns_actual * 100) if ns_actual else 0
        ly_pct_ns = (ly / ns_ly * 100) if ns_ly else 0
        delta = pct_ns - ly_pct_ns

        cell = table.cell(row_idx, 0)
        cell.text = display_name
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True

        _style_data_cell(table.cell(row_idx, 1), actual)
        _style_data_cell(table.cell(row_idx, 2), aop)
        # For costs, positive variance = bad
        _style_data_cell(table.cell(row_idx, 3), var_pct, False, False,
                         RED if var_pct > 5 else GREEN if var_pct < -5 else BLACK)
        _style_data_cell(table.cell(row_idx, 4), pct_ns, False)
        _style_data_cell(table.cell(row_idx, 5), ly_pct_ns, False)
        _style_data_cell(table.cell(row_idx, 6), delta, False, False,
                         RED if delta > 0.5 else GREEN if delta < -0.5 else BLACK)

    # ==================== SLIDE 13: Outlier Summary ====================
    slide = _add_slide(prs)
    _add_title_bar(slide, "Outlier Summary — All Comparisons",
                   f"Total {len(a.outliers)} outliers detected across AOP, MoM, YoY, QoQ")

    # Count by type and severity
    type_counts = {}
    severity_counts = {"high": 0, "medium": 0}
    for o in a.outliers:
        type_counts[o.comparison_type] = type_counts.get(o.comparison_type, 0) + 1
        severity_counts[o.severity] += 1

    _add_text_box(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.4),
                  "OUTLIER DISTRIBUTION", font_size=12, bold=True, color=DARK_BLUE)

    y = 1.8
    for typ, count in type_counts.items():
        label = {"vs_AOP": "vs Budget (AOP)", "MoM": "Month-over-Month",
                 "YoY": "Year-over-Year", "QoQ": "Quarter-over-Quarter"}.get(typ, typ)
        _add_text_box(slide, Inches(0.7), Inches(y), Inches(6), Inches(0.3),
                      f"{label}: {count} outliers", font_size=10)
        y += 0.35

    _add_text_box(slide, Inches(0.5), Inches(y + 0.3), Inches(12), Inches(0.4),
                  f"HIGH SEVERITY ({severity_counts['high']})", font_size=12, bold=True, color=RED)

    y += 0.8
    high_outliers = [o for o in a.outliers if o.severity == "high"][:10]
    for o in high_outliers:
        _add_text_box(slide, Inches(0.7), Inches(y), Inches(12), Inches(0.3),
                      f"[{o.comparison_type}] {o.description}", font_size=8, color=RED)
        y += 0.3

    if severity_counts["medium"] > 0:
        _add_text_box(slide, Inches(0.5), Inches(y + 0.2), Inches(12), Inches(0.4),
                      f"MEDIUM SEVERITY ({severity_counts['medium']})", font_size=12, bold=True, color=ORANGE)
        y += 0.6
        med_outliers = [o for o in a.outliers if o.severity == "medium"][:8]
        for o in med_outliers:
            _add_text_box(slide, Inches(0.7), Inches(y), Inches(12), Inches(0.3),
                          f"[{o.comparison_type}] {o.description}", font_size=8, color=ORANGE)
            y += 0.3

    # ==================== LLM INSIGHT SLIDES (if available) ====================
    if insights:
        # Executive Summary Narrative
        if "executive_summary" in insights:
            _add_narrative_slide(prs,
                                f"AI-Powered Executive Summary — {period_label}",
                                "Generated by Claude | Key takeaways for leadership",
                                insights["executive_summary"])

        # Budget Commentary
        if "budget_commentary" in insights:
            bullets = insights["budget_commentary"]
            if isinstance(bullets, str):
                bullets = [s.strip() for s in bullets.split(".") if s.strip()]
            _add_narrative_slide(prs,
                                "Budget Achievement Commentary",
                                "AI analysis of AOP performance",
                                bullets)

        # Cost Commentary
        if "cost_commentary" in insights:
            bullets = insights["cost_commentary"]
            if isinstance(bullets, str):
                bullets = [s.strip() for s in bullets.split(".") if s.strip()]
            _add_narrative_slide(prs,
                                "Cost Trend Analysis",
                                "AI analysis of cost structure and trends",
                                bullets)

        # Outlier Commentary
        if "outlier_commentary" in insights:
            _add_narrative_slide(prs,
                                "Outlier Analysis — AI Commentary",
                                "AI-generated explanations for significant deviations",
                                insights["outlier_commentary"])

        # Risks & Actions
        if "risks_and_actions" in insights:
            _add_narrative_slide(prs,
                                "Risks & Recommended Actions",
                                "AI-identified risks and suggested next steps",
                                insights["risks_and_actions"], RED)

        # Quarter Commentary
        if "quarter_commentary" in insights:
            bullets = insights["quarter_commentary"]
            if isinstance(bullets, str):
                bullets = [s.strip() for s in bullets.split(".") if s.strip()]
            _add_narrative_slide(prs,
                                "Quarterly Performance Commentary",
                                "AI analysis of quarterly trends",
                                bullets)

        # Full Year (if available)
        if "full_year_commentary" in insights and is_full_year:
            bullets = insights["full_year_commentary"]
            if isinstance(bullets, str):
                bullets = [s.strip() for s in bullets.split(".") if s.strip()]
            _add_narrative_slide(prs,
                                f"Full Year Commentary — {a.review_fy}",
                                "AI analysis of annual performance",
                                bullets)

    # Save
    prs.save(output_path)
    return output_path


def _build_outlier_table(slide, outliers, comparison_label):
    """Build a table of outliers on a slide."""
    cols = ["Line Item", "Actual", "Reference", "Deviation %", "Direction", "Severity"]
    num_rows = len(outliers) + 1
    table = _add_table(slide, Inches(0.3), Inches(1.1), Inches(12.7),
                       Inches(min(6.2, 0.4 + len(outliers) * 0.35)),
                       num_rows, len(cols))

    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(1.8)
    table.columns[3].width = Inches(1.8)
    table.columns[4].width = Inches(1.8)
    table.columns[5].width = Inches(2.0)

    for ci, h in enumerate(cols):
        table.cell(0, ci).text = h
    _style_header_row(table, len(cols))

    for ri, o in enumerate(outliers):
        row_idx = ri + 1

        cell = table.cell(row_idx, 0)
        cell.text = o.line_item
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.bold = True

        _style_data_cell(table.cell(row_idx, 1), o.current_value)
        _style_data_cell(table.cell(row_idx, 2), o.reference_value)

        dev_color = RED if o.direction == "above" else GREEN
        # For revenue items, above = good
        if any(kw in o.line_item for kw in ["Sales", "Revenue", "Margin", "Profit", "EBITDA", "EBIT", "PAT", "Income"]):
            dev_color = GREEN if o.direction == "above" else RED

        _style_data_cell(table.cell(row_idx, 3), o.deviation_pct, False, False, dev_color)

        cell = table.cell(row_idx, 4)
        cell.text = o.direction.upper()
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.color.rgb = dev_color
            p.alignment = PP_ALIGN.CENTER

        cell = table.cell(row_idx, 5)
        cell.text = o.severity.upper()
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = RED if o.severity == "high" else ORANGE
            p.alignment = PP_ALIGN.CENTER

        if o.severity == "high":
            table.cell(row_idx, 5).fill.solid()
            table.cell(row_idx, 5).fill.fore_color.rgb = RGBColor(0xFF, 0xE0, 0xE0)
